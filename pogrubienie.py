from pptx import Presentation
import sys
import re

def bold_labels_before_colon(input_path, output_path):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = paragraph.text
                    paragraph.clear()

                    # Szukamy wzorców: wielka litera, coś-tam, dwukropek
                    matches = list(re.finditer(r'[A-ZĄĆĘŁŃÓŚŹŻ][^:]{1,50}:', full_text))

                    cursor = 0
                    for match in matches:
                        start, end = match.span()
                        if cursor < start:
                            run = paragraph.add_run()
                            run.text = full_text[cursor:start]
                        run = paragraph.add_run()
                        run.text = full_text[start:end]
                        run.font.bold = True
                        cursor = end

                    # Dodaj resztę tekstu
                    if cursor < len(full_text):
                        run = paragraph.add_run()
                        run.text = full_text[cursor:]

    prs.save(output_path)
    print(f"✅ Zapisano: {output_path}")

# Obsługa terminala
if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Użycie: python pogrub.py prezentacja.pptx nowa_prezentacja.pptx")
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        bold_labels_before_colon(input_file, output_file)
