from pptx import Presentation
from deep_translator import GoogleTranslator

translator = GoogleTranslator(source='en', target='pl')

def translate_pptx(input_path, output_path):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines = shape.text_frame.text.split("\n")
                shape.text_frame.clear()
                for idx, line in enumerate(lines):
                    if line.strip():
                        translated = translator.translate(line)
                    else:
                        translated = ""
                    if idx == 0:
                        shape.text_frame.text = translated
                    else:
                        shape.text_frame.add_paragraph().text = translated

    prs.save(output_path)
    print(f"Przetłumaczono: {output_path}")

# Przykład użycia:
translate_pptx("part1.pptx", "avaPL.pptx")
