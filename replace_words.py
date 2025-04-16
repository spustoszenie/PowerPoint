import sys
from pptx import Presentation

replacements = {
    "Cele uczenia się": "Cele nauki",
    "Ćwicz zadania przygotowanie": "Przygotowanie do zadań praktycznych",
    "Zadanie ćwiczeń":"Ćwiczenie praktyczne",
    "Ćwicz rozwiązanie zadań":"Rozwiązanie",
    "Samocena":"Praca własna",
    "Strategie łagodzenia":"Środki zaradcze"

}

def replace_text_in_pptx(input_path, output_path):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for old, new in replacements.items():
                        if old in paragraph.text:
                            paragraph.text = paragraph.text.replace(old, new)

    prs.save(output_path)
    print(f"✅ Zapisano zmodyfikowaną prezentację jako: {output_path}")

# Obsługa z terminala
if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Użycie: python zamiana.py prezentacja.pptx nowa_prezentacja.pptx")
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        replace_text_in_pptx(input_file, output_file)
